#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
SynchroniNote 紹介スライド: Markdown -> PowerPoint(.pptx) 変換

- 内容は同じフォルダの `SynchroniNote_紹介スライド.md`（素のMarkdown）で編集する。
- 本スクリプトを実行すると `SynchroniNote_紹介スライド.pptx` を書き出す（提示用）。
- 依存: python-pptx（`pip install python-pptx` / 本環境は導入済み）。

Markdown の書式ルール（このスクリプトが解釈する範囲）:
  `# 見出し`            … 1スライドの開始（タイトル）。最初の `#` はタイトルスライド。
  `## サブタイトル`     … タイトルスライドのサブタイトル。
  `- 箇条書き`          … 箇条書き（先頭2スペース字下げで小項目）。
  `> ひとこと`          … 強調コールアウト（下部の黄色い帯）。
  `![説明](images/x.png)` … 画像（1枚=右、2枚=下段2列）。
  `<!-- note: ... -->`  … 発表者ノート（PowerPointのノート欄へ）。複数行可。
  それ以外の行          … 箇条書きの上の導入文。

使い方:  python build_pptx.py
"""
import os
import re
import sys
import struct

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

# ---- テーマ ----------------------------------------------------------------
NAVY        = RGBColor(0x1E, 0x3A, 0x8A)
GOLD        = RGBColor(0xFC, 0xD3, 0x4D)
AMBER_TEXT  = RGBColor(0x92, 0x40, 0x0E)
AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
GRAY        = RGBColor(0x1F, 0x29, 0x37)
LIGHTGRAY   = RGBColor(0x64, 0x74, 0x8B)
SILVER      = RGBColor(0xCB, 0xD5, 0xE1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
FONT        = "Meiryo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TITLE_H = Inches(1.05)

HERE = os.path.dirname(os.path.abspath(__file__))


# ---- 低レベルヘルパ --------------------------------------------------------
def set_run(run, text, size, bold, color):
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = FONT  # latin
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    prev = latin
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            if prev is not None:
                prev.addnext(el)
            else:
                rPr.append(el)
        el.set('typeface', FONT)
        prev = el


def no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def png_size(path):
    with open(path, 'rb') as f:
        head = f.read(33)
    if head[:8] != b'\x89PNG\r\n\x1a\n':
        return (1800, 1200)
    w, h = struct.unpack('>II', head[16:24])
    return (w, h)


def resolve(p):
    return p if os.path.isabs(p) else os.path.join(HERE, p)


def read_dd_entries():
    """doc/DD/DD-INDEX.md から (DD番号, 件名) を抽出。件名は長すぎたら truncate。"""
    path = os.path.join(HERE, '..', 'DD', 'DD-INDEX.md')
    entries = []
    if not os.path.exists(path):
        return entries
    with open(path, encoding='utf-8') as f:
        for line in f:
            m = re.match(r'\|\s*(DD-[0-9A-Za-z-]+)\s*\|\s*([^|]+?)\s*\|', line)
            if m:
                title = m.group(2).strip().replace('_', ' ')
                if len(title) > 16:
                    title = title[:16] + '…'
                entries.append((m.group(1).strip(), title))
    return entries


# ---- 部品 ------------------------------------------------------------------
def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, TITLE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background(); no_shadow(bar)
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), TITLE_H)
    acc.fill.solid(); acc.fill.fore_color.rgb = GOLD; acc.line.fill.background(); no_shadow(acc)
    tb = slide.shapes.add_textbox(Inches(0.45), 0, SLIDE_W - Inches(0.9), TITLE_H)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    set_run(p.add_run(), title, 25, True, WHITE)


def add_chapter_rail(slide, chapters, active):
    """見出し帯の下に章バー（ステッパー）を置く。現在の章を金色で強調。"""
    n = len(chapters)
    if n == 0:
        return
    M = Inches(0.45)
    top = Inches(1.14)
    h = Inches(0.4)
    gap = Inches(0.1)
    pw = int((SLIDE_W - 2 * M - gap * (n - 1)) / n)
    x = int(M)
    for name in chapters:
        is_active = (name == active)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, pw, h)
        if is_active:
            pill.fill.solid(); pill.fill.fore_color.rgb = GOLD; pill.line.fill.background()
            tcolor, bold = NAVY, True
        else:
            pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xF7)
            pill.line.color.rgb = SILVER; pill.line.width = Pt(0.75)
            tcolor, bold = LIGHTGRAY, False
        no_shadow(pill)
        tf = pill.text_frame; tf.word_wrap = False
        tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
        tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), name, 11, bold, tcolor)
        x = int(x + pw + gap)


def add_page_number(slide, n, total):
    tb = slide.shapes.add_textbox(SLIDE_W - Inches(1.35), SLIDE_H - Inches(0.36), Inches(1.15), Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    set_run(p.add_run(), '%d / %d' % (n, total), 11, False, LIGHTGRAY)


def add_paras(slide, paras, left, top, width, size=15, color=LIGHTGRAY):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        add_rich(p, pa, size, color, NAVY)
    return tb


def _is_letter_start(s):
    if not s:
        return False
    import unicodedata
    return unicodedata.category(s[0])[0] in ('L', 'N')


def add_rich(paragraph, text, size, base_color, emph_color):
    """`**...**` を太字（emph_color）に変換しながらランを追加する。"""
    for i, seg in enumerate(text.split('**')):
        if seg == '':
            continue
        emph = (i % 2 == 1)
        set_run(paragraph.add_run(), seg, size, emph, emph_color if emph else base_color)


def add_bullets(slide, bullets, left, top, width, height, size=18):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(9)
        if b['level'] == 1:
            p.level = 1
            add_rich(p, '– ' + b['text'], max(size - 4, 12), LIGHTGRAY, NAVY)
        else:
            if _is_letter_start(b['text']):
                set_run(p.add_run(), '•  ', size, False, GRAY)
            add_rich(p, b['text'], size, GRAY, NAVY)
    return tb


def add_callout(slide, text, left, top, width, height):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = AMBER_LIGHT
    box.line.color.rgb = GOLD; box.line.width = Pt(1.25); no_shadow(box)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text.replace('**', ''), 15, True, AMBER_TEXT)
    return box


def add_image_fit(slide, path, box_l, box_t, box_w, box_h, border=True):
    if not os.path.exists(path):
        ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_l, box_t, box_w, box_h)
        ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
        ph.line.color.rgb = SILVER; no_shadow(ph)
        tf = ph.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), '［画像なし: %s］' % os.path.basename(path), 12, False, LIGHTGRAY)
        return ph
    w, h = png_size(path)
    ar = w / h
    if ar > (box_w / box_h):
        dw = box_w; dh = int(box_w / ar)
    else:
        dh = box_h; dw = int(box_h * ar)
    left = int(box_l + (box_w - dw) / 2)
    top = int(box_t + (box_h - dh) / 2)
    pic = slide.shapes.add_picture(path, left, top, width=dw, height=dh)
    if border:
        pic.line.color.rgb = SILVER; pic.line.width = Pt(1)
    return pic


def add_image(slide, path, alt, box_l, box_t, box_w, box_h):
    cap_h = Inches(0.32)
    add_image_fit(slide, resolve(path), box_l, box_t, box_w, box_h - cap_h)
    tb = slide.shapes.add_textbox(box_l, box_t + box_h - cap_h, box_w, cap_h)
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), alt, 10, False, LIGHTGRAY)


# ---- スライド描画 ----------------------------------------------------------
def render_title(slide, s):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); no_shadow(bg)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.25), SLIDE_W - Inches(2.0), Inches(1.4))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), s['title'], 54, True, WHITE)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(SLIDE_W / 2 - Inches(1.25)), Inches(3.75), Inches(2.5), Pt(3))
    rule.fill.solid(); rule.fill.fore_color.rgb = GOLD; rule.line.fill.background(); no_shadow(rule)
    if s['subtitle']:
        sb = slide.shapes.add_textbox(Inches(1.0), Inches(3.95), SLIDE_W - Inches(2.0), Inches(0.9))
        tf = sb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), s['subtitle'], 24, False, GOLD)
    if s['paras']:
        fb = slide.shapes.add_textbox(Inches(1.0), Inches(5.1), SLIDE_W - Inches(2.0), Inches(1.4))
        tf = fb.text_frame; tf.word_wrap = True
        for i, pa in enumerate(s['paras']):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
            set_run(p.add_run(), pa, 16, False, SILVER)


def render_content(slide, s, chapters, page_no, total):
    add_title_bar(slide, s['title'])
    if s.get('chapter') and s['chapter'] in chapters:
        add_chapter_rail(slide, chapters, s['chapter'])
        body_top = Inches(1.66)
    else:
        body_top = Inches(1.32)
    add_page_number(slide, page_no, total)

    imgs = s['images']
    intro = s['paras']
    callout = ' '.join(s['callout']) if s['callout'] else None
    M = Inches(0.6)

    if s.get('dd_table'):
        entries = read_dd_entries()
        cur = body_top
        head = '全 %d 件 — これらすべての「判断」が設計ドキュメントとして残る（監査でき・引き継げる）' % len(entries)
        add_paras(slide, [head], M, cur, SLIDE_W - 2 * M, size=13, color=GRAY)
        cur += Inches(0.45)
        ncol = 3
        gap = Inches(0.3)
        colw = int((SLIDE_W - 2 * M - gap * (ncol - 1)) / ncol)
        per = (len(entries) + ncol - 1) // ncol
        colh = int(SLIDE_H - cur - Inches(0.4))
        for ci in range(ncol):
            chunk = entries[ci * per:(ci + 1) * per]
            x = int(M + ci * (colw + gap))
            tb = slide.shapes.add_textbox(x, cur, colw, colh)
            tf = tb.text_frame; tf.word_wrap = True
            for j, (eid, etitle) in enumerate(chunk):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_after = Pt(2)
                set_run(p.add_run(), eid + ' ', 10, True, NAVY)
                set_run(p.add_run(), etitle, 10, False, GRAY)
        return

    if s.get('wide') and imgs:
        cur = body_top
        if intro:
            add_paras(slide, intro, M, cur, SLIDE_W - 2 * M, size=15, color=GRAY)
            cur += Inches(0.42)
        if s['bullets']:
            bh = Inches(1.3)
            add_bullets(slide, s['bullets'], M, cur, SLIDE_W - 2 * M, bh, size=15)
            cur += bh + Inches(0.08)
        add_image(slide, imgs[0]['path'], imgs[0]['alt'], M, cur,
                  SLIDE_W - 2 * M, int(SLIDE_H - cur - Inches(0.45)))
        return

    if len(imgs) >= 2:
        cur = body_top
        if intro:
            add_paras(slide, intro, M, cur, SLIDE_W - 2 * M); cur += Inches(0.46)
        bh = Inches(1.65)
        add_bullets(slide, s['bullets'], M, cur, SLIDE_W - 2 * M, bh, size=16)
        band_t = int(cur + bh + Inches(0.1))
        band_h = int(SLIDE_H - band_t - Inches(0.5))
        gap = Inches(0.35)
        half = int((SLIDE_W - 2 * M - gap) / 2)
        add_image(slide, imgs[0]['path'], imgs[0]['alt'], M, band_t, half, band_h)
        add_image(slide, imgs[1]['path'], imgs[1]['alt'], int(M + half + gap), band_t, half, band_h)
    elif len(imgs) == 1:
        left_w = Inches(6.4)
        cur = body_top
        if intro:
            add_paras(slide, intro, M, cur, left_w); cur += Inches(0.5)
        add_bullets(slide, s['bullets'], M, cur, left_w, Inches(3.8), size=18)
        if callout:
            add_callout(slide, callout, M, Inches(6.35), left_w, Inches(0.8))
        add_image(slide, imgs[0]['path'], imgs[0]['alt'], Inches(7.25), body_top, Inches(5.55), Inches(4.7))
    else:
        cur = body_top
        if intro:
            add_paras(slide, intro, Inches(0.75), cur, SLIDE_W - Inches(1.5), size=17, color=GRAY)
            cur += Inches(0.55)
        bh = Inches(4.2) if callout else Inches(5.1)
        add_bullets(slide, s['bullets'], Inches(0.85), cur, SLIDE_W - Inches(1.7), bh, size=19)
        if callout:
            add_callout(slide, callout, Inches(0.85), Inches(6.28), SLIDE_W - Inches(1.7), Inches(0.8))


# ---- Markdown パーサ -------------------------------------------------------
def parse(md):
    slides = []
    chapters = []
    cur = None
    current_chapter = None
    pending_wide = False
    pending_dd = False
    in_note = False
    note_buf = []
    for line in md.splitlines():
        if in_note:
            if '-->' in line:
                note_buf.append(line.split('-->')[0])
                if cur is not None:
                    cur['notes'] = '\n'.join(note_buf).strip()
                in_note = False; note_buf = []
            else:
                note_buf.append(line)
            continue
        st = line.strip()
        if st.startswith('<!-- chapter:'):
            name = st[len('<!-- chapter:'):].split('-->')[0].strip()
            current_chapter = name or None
            if name and name not in chapters:
                chapters.append(name)
            continue
        if st.startswith('<!-- appendix'):
            current_chapter = None
            continue
        if st.startswith('<!-- wide'):
            pending_wide = True
            continue
        if st.startswith('<!-- dd-table'):
            pending_dd = True
            continue
        if st.startswith('<!-- note:'):
            content = st[len('<!-- note:'):]
            if '-->' in content:
                if cur is not None:
                    cur['notes'] = content.split('-->')[0].strip()
            else:
                in_note = True; note_buf = [content]
            continue
        if st.startswith('# '):
            cur = {'title': st[2:].strip(), 'subtitle': None, 'bullets': [],
                   'callout': [], 'images': [], 'paras': [], 'notes': '',
                   'chapter': current_chapter, 'wide': pending_wide,
                   'dd_table': pending_dd}
            pending_wide = False
            pending_dd = False
            slides.append(cur); continue
        if cur is None:
            continue
        if st.startswith('## '):
            cur['subtitle'] = st[3:].strip()
        elif st.startswith('!['):
            m = re.match(r'!\[(.*?)\]\((.*?)\)', st)
            if m:
                cur['images'].append({'alt': m.group(1), 'path': m.group(2)})
        elif st.startswith('> '):
            cur['callout'].append(st[2:].strip())
        elif st.startswith('- '):
            indent = len(line) - len(line.lstrip())
            cur['bullets'].append({'level': 1 if indent >= 2 else 0, 'text': st[2:].strip()})
        elif st == '':
            pass
        else:
            cur['paras'].append(st)
    return slides, chapters


def main():
    md_path = os.path.join(HERE, 'SynchroniNote_紹介スライド.md')
    out_path = os.path.join(HERE, 'SynchroniNote_紹介スライド.pptx')
    with open(md_path, encoding='utf-8') as f:
        md = f.read()
    slides, chapters = parse(md)
    total = len(slides)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    for idx, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        if idx == 0:
            render_title(slide, s)
        else:
            render_content(slide, s, chapters, idx + 1, total)
        if s['notes']:
            slide.notes_slide.notes_text_frame.text = s['notes']
    prs.save(out_path)
    print('Saved: %s  (%d slides, chapters: %s)' % (out_path, len(slides), ' / '.join(chapters)))


if __name__ == '__main__':
    main()
